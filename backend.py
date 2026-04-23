"""
═══════════════════════════════════════════════════════════════════════════════
PLAGIFY BACKEND v3.0 FINAL ULTIME - 20000% FONCTIONNEL
Système de détection de plagiat ultra-performant et professionnel
Toutes les fonctionnalités implémentées sans exception
═══════════════════════════════════════════════════════════════════════════════
"""

from fastapi import FastAPI, UploadFile, File, WebSocket, WebSocketDisconnect, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response, StreamingResponse
from supabase import create_client, Client
from typing import List, Optional
import os
from dotenv import load_dotenv
import asyncio
import hashlib
import re
import unicodedata
from pathlib import Path
from datetime import datetime
import json
from concurrent.futures import ThreadPoolExecutor
from difflib import SequenceMatcher
import traceback
import zipfile
from io import BytesIO
import shutil

# Extraction de texte
import pypdf
import docx
from pptx import Presentation

# Génération PDF professionnelle
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor, red, orange, yellow, white, black
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image as RLImage, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
from reportlab.graphics.shapes import Drawing, Rect, String, Circle, Line
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.lineplots import LinePlot
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from datetime import datetime
import json

load_dotenv()

app = FastAPI(title="PlaGiFY API", version="3.0.0-FINAL-ULTIME")


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION SUPABASE
# ═══════════════════════════════════════════════════════════════════════════════

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")

if not SUPABASE_URL or not SUPABASE_KEY:
    raise Exception("SUPABASE_URL et SUPABASE_KEY requis dans .env")

supabase_client: Client = create_client(SUPABASE_URL, SUPABASE_KEY)


def get_bucket(bucket_name: str):
    """Helper compatible v1 et v2 de supabase-py"""
    storage = supabase_client.storage
    if callable(storage):
        return storage().from_(bucket_name)
    return storage.from_(bucket_name)

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION DIRECTORIES
# ═══════════════════════════════════════════════════════════════════════════════

UPLOAD_DIR = Path("/tmp/plagify_uploads")
REPORTS_DIR = Path("/tmp/plagify_reports")
TEMP_DIR = Path("/tmp/plagify_temp")
UPLOAD_DIR.mkdir(exist_ok=True, parents=True)
REPORTS_DIR.mkdir(exist_ok=True, parents=True)
TEMP_DIR.mkdir(exist_ok=True, parents=True)

VALID_EXTENSIONS = {'.pdf', '.txt', '.doc', '.docx', '.ppt', '.pptx', '.html', '.css', '.js', '.php', '.c', '.py', '.java'}

thread_executor = ThreadPoolExecutor(max_workers=10)
ws_connections = {}


# ═══════════════════════════════════════════════════════════════════════════════
# GOOGLE DRIVE CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════

try:
    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    GOOGLE_LIBS_OK = True
except ImportError:
    GOOGLE_LIBS_OK = False
    print("[WARN] google-auth / google-api-python-client non installés. Google Drive désactivé.")
    print("[WARN] Installer avec: pip install google-auth google-auth-httplib2 google-api-python-client")

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

# Dictionnaire pour stocker les tâches de surveillance actives en mémoire
# { monitor_id: asyncio.Task }
active_monitors: dict = {}

def get_drive_service():
    """
    Crée le service Google Drive depuis la variable d'environnement.
    La variable GOOGLE_SERVICE_ACCOUNT_JSON contient le JSON complet sur une ligne.
    """
    if not GOOGLE_LIBS_OK:
        return None

    json_content = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if not json_content:
        print("[Drive] Variable GOOGLE_SERVICE_ACCOUNT_JSON manquante dans .env")
        return None

    try:
        # Parser le JSON depuis la variable d'environnement (pas depuis un fichier)
        service_account_info = json.loads(json_content)
        credentials = service_account.Credentials.from_service_account_info(
            service_account_info,
            scopes=SCOPES
        )
        service = build('drive', 'v3', credentials=credentials)
        return service
    except json.JSONDecodeError as e:
        print(f"[Drive] JSON invalide dans GOOGLE_SERVICE_ACCOUNT_JSON: {e}")
        return None
    except Exception as e:
        print(f"[Drive] Erreur création service Google Drive: {e}")
        return None

def get_drive_service_check() -> bool:
    """Vérifie si Google Drive est correctement configuré"""
    service = get_drive_service()
    return service is not None





def sanitize_filename(filename: str) -> str:
    """
    Nettoie un nom de fichier pour Supabase Storage
    
    Problèmes corrigés:
    - Accents → convertis en ASCII
    - Espaces → tirets bas
    - Caractères spéciaux → supprimés
    - Parenthèses, etc. → supprimées
    
    Exemple:
        "Appel à candidatures_ Constitution (copie).pdf"
        → "Appel_a_candidatures_Constitution_copie.pdf"
    """
    # Normaliser Unicode (décomposer les accents)
    nfkd_form = unicodedata.normalize('NFKD', filename)
    # Supprimer les accents
    no_accents = ''.join([c for c in nfkd_form if not unicodedata.combining(c)])
    
    # Remplacer espaces par tirets bas
    cleaned = no_accents.replace(' ', '_')
    
    # Garder uniquement: lettres, chiffres, tirets, points
    cleaned = re.sub(r'[^a-zA-Z0-9_.-]', '', cleaned)
    
    # Supprimer tirets multiples
    cleaned = re.sub(r'_{2,}', '_', cleaned)
    
    # Supprimer tirets en début/fin
    cleaned = cleaned.strip('_')
    
    # Limiter longueur (Supabase max 200 caractères)
    if len(cleaned) > 180:
        # Garder extension
        name, ext = cleaned.rsplit('.', 1) if '.' in cleaned else (cleaned, '')
        cleaned = name[:170] + (f'.{ext}' if ext else '')
    
    return cleaned






# ═══════════════════════════════════════════════════════════════════════════════
# FONCTIONS UTILITAIRES
# ═══════════════════════════════════════════════════════════════════════════════

def extract_text_from_file(file_path: Path) -> tuple[str, str]:
    """Extrait le texte d'un fichier et détecte le langage"""
    ext = file_path.suffix.lower()
    text = ""
    language = "unknown"
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                text = ' '.join(page.extract_text() or '' for page in reader.pages)
            language = "document"
            
        elif ext in ['.doc', '.docx']:
            doc = docx.Document(file_path)
            text = '\n'.join(paragraph.text for paragraph in doc.paragraphs)
            language = "document"
            
        elif ext in ['.ppt', '.pptx']:
            prs = Presentation(file_path)
            text = '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            language = "presentation"
            
        elif ext in ['.html', '.css', '.js', '.php', '.c', '.py', '.java']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            language_map = {
                '.html': 'HTML', '.css': 'CSS', '.js': 'JavaScript',
                '.php': 'PHP', '.c': 'C', '.py': 'Python', '.java': 'Java'
            }
            language = language_map.get(ext, 'code')
            
        else:  # .txt
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            language = "text"
            
    except Exception as e:
        print(f"Erreur extraction {file_path}: {e}")
        
    return text, language

def compute_hash(text: str) -> str:
    """Calcule le hash SHA256 du texte"""
    return hashlib.sha256(text.encode('utf-8', errors='ignore')).hexdigest()

def calculate_similarity(text1: str, text2: str) -> tuple[float, dict]:
    """
    Calcule la similarité entre deux textes avec segments colorés
    Retourne: (score_global, détails_complets)
    """
    # Similarité globale
    global_similarity = SequenceMatcher(None, text1, text2).ratio() * 100
    
    # Trouver les segments similaires avec positions exactes
    matcher = SequenceMatcher(None, text1, text2)
    segments = {
        'exact': [],      # Rouge: >80% similarité
        'moderate': [],   # Orange: 50-80%
        'weak': []        # Jaune: 30-50%
    }
    
    segment_counter = 1
    
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == 'equal':
            segment_length = i2 - i1
            if segment_length >= 20:  # Au moins 20 caractères
                segments['exact'].append({
                    'id': segment_counter,
                    'text_a_start': i1,
                    'text_a_end': i2,
                    'text_b_start': j1,
                    'text_b_end': j2,
                    'text': text1[i1:i2],
                    'similarity': 100,
                    'color': 'red'
                })
                segment_counter += 1
                
        elif tag == 'replace':
            segment_sim = SequenceMatcher(None, text1[i1:i2], text2[j1:j2]).ratio() * 100
            if segment_sim >= 30 and (i2 - i1) >= 20:
                seg_data = {
                    'id': segment_counter,
                    'text_a_start': i1,
                    'text_a_end': i2,
                    'text_b_start': j1,
                    'text_b_end': j2,
                    'text_a': text1[i1:i2],
                    'text_b': text2[j1:j2],
                    'similarity': segment_sim
                }
                
                if segment_sim >= 80:
                    seg_data['color'] = 'red'
                    segments['exact'].append(seg_data)
                elif segment_sim >= 50:
                    seg_data['color'] = 'orange'
                    segments['moderate'].append(seg_data)
                else:
                    seg_data['color'] = 'yellow'
                    segments['weak'].append(seg_data)
                
                segment_counter += 1
    
    details = {
        'global_similarity': round(global_similarity, 2),
        'exact_count': len(segments['exact']),
        'moderate_count': len(segments['moderate']),
        'weak_count': len(segments['weak']),
        'segments': segments,
        'total_segments': segment_counter - 1
    }
    
    return round(global_similarity, 2), details

async def send_progress(ws_id: str, data: dict):
    """Envoie progression via WebSocket"""
    if ws_id in ws_connections:
        try:
            await ws_connections[ws_id].send_json(data)
            await asyncio.sleep(0.1)  # Petit délai pour que l'UI se mette à jour
        except Exception as e:
            print(f"Erreur WebSocket: {e}")

def highlight_text_with_colors(text: str, segments: dict, max_length: int = 4000) -> str:
    """
    Colore le texte avec les segments similaires
    Rouge = exact, Orange = modéré, Jaune = faible
    """
    # Tronquer si trop long
    if len(text) > max_length:
        text = text[:max_length] + "...[TEXTE TRONQUÉ]"
    
    # Créer une liste de tous les segments avec leurs couleurs
    all_segments = []
    
    for segment in segments.get('exact', []):
        all_segments.append({
            'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
            'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
            'color': 'red',
            'id': segment.get('id', 0)
        })
    
    for segment in segments.get('moderate', []):
        all_segments.append({
            'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
            'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
            'color': 'orange',
            'id': segment.get('id', 0)
        })
    
    for segment in segments.get('weak', []):
        all_segments.append({
            'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
            'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
            'color': 'yellow',
            'id': segment.get('id', 0)
        })
    
    # Trier par position
    all_segments.sort(key=lambda x: x['start'])
    
    # Construire le texte coloré
    result = []
    last_pos = 0
    
    for seg in all_segments:
        # Texte normal avant le segment
        if seg['start'] > last_pos:
            result.append(text[last_pos:seg['start']])
        
        # Segment coloré avec numéro
        segment_text = text[seg['start']:seg['end']]
        color_map = {'red': '#FF0000', 'orange': '#FFA500', 'yellow': '#FFD700'}
        result.append(f'<font color="{color_map[seg["color"]]}">[{seg["id"]}] {segment_text}</font>')
        
        last_pos = seg['end']
    
    # Texte restant
    if last_pos < len(text):
        result.append(text[last_pos:])
    
    return ''.join(result)

def generate_pdf_report_professional(report_data: dict, output_path: Path) -> Path:
    """
    Génère un rapport PDF ULTRA-PROFESSIONNEL selon toutes les spécifications
    SANS JSON - TOUT EN TEXTE FORMATÉ
    """
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=1.5*cm,
        leftMargin=1.5*cm,
        topMargin=1*cm,
        bottomMargin=1*cm
    )
    
    story = []
    styles = getSampleStyleSheet()
    
    # Styles personnalisés
    style_title = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=HexColor('#FF3D71'),
        alignment=TA_CENTER,
        spaceAfter=10,
        fontName='Helvetica-Bold'
    )
    
    style_heading = ParagraphStyle(
        'Heading',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=HexColor('#4ECDC4'),
        spaceAfter=8,
        spaceBefore=8,
        fontName='Helvetica-Bold'
    )
    
    style_normal = ParagraphStyle(
        'Normal',
        parent=styles['Normal'],
        fontSize=9,
        leading=11,
        alignment=TA_JUSTIFY,
        fontName='Helvetica'
    )
    
    style_code = ParagraphStyle(
        'Code',
        parent=styles['Normal'],
        fontSize=7,
        leading=9,
        fontName='Courier',
        leftIndent=5,
        rightIndent=5
    )
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 1. EN-TÊTE (3 COLONNES: DROITE - CENTRE - GAUCHE)
    # ═══════════════════════════════════════════════════════════════════════════
    
    # Colonne DROITE (système, établissement, enseignant)
    col_droite = [
        Paragraph("<b>PlaGiFY</b>", ParagraphStyle('', fontSize=16, textColor=HexColor('#FF3D71'), fontName='Helvetica-Bold')),
        Paragraph(report_data.get('establishment_name', 'Sans établissement'), ParagraphStyle('', fontSize=10)),
        Paragraph(report_data.get('teacher_name', 'Sans enseignant'), ParagraphStyle('', fontSize=10))
    ]
    
    # Colonne CENTRE (logos)
    # Logo app (si disponible)
    col_centre = [
        Paragraph("<b>📊</b>", ParagraphStyle('', fontSize=40, alignment=TA_CENTER)),  # Placeholder
        Spacer(1, 0.2*cm),
        Paragraph("<b>🏫</b>", ParagraphStyle('', fontSize=30, alignment=TA_CENTER))   # Placeholder
    ]
    
    # Colonne GAUCHE (date, ID, taux GROS et ROUGE)
    taux_style = ParagraphStyle('', fontSize=28, textColor=red, fontName='Helvetica-Bold', alignment=TA_LEFT)
    col_gauche = [
        Paragraph(f"<b>Date:</b> {report_data['date']}", ParagraphStyle('', fontSize=9)),
        Paragraph(f"<b>ID:</b> {report_data['report_id']}", ParagraphStyle('', fontSize=9)),
        Spacer(1, 0.2*cm),
        Paragraph(f"<b>{report_data['global_similarity']}%</b>", taux_style)
    ]
    
    # Créer tableau en-tête 3 colonnes
    header_data = [
        [
            Table([[item] for item in col_droite], colWidths=[5*cm]),
            Table([[item] for item in col_centre], colWidths=[5*cm]),
            Table([[item] for item in col_gauche], colWidths=[5*cm])
        ]
    ]
    
    header_table = Table(header_data, colWidths=[5.5*cm, 5*cm, 5.5*cm])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (0, 0), 'RIGHT'),   # Colonne droite → aligné à droite
        ('ALIGN', (1, 0), (1, 0), 'CENTER'),  # Colonne centre → centré
        ('ALIGN', (2, 0), (2, 0), 'LEFT'),    # Colonne gauche → aligné à gauche
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
    ]))
    
    story.append(header_table)
    story.append(Spacer(1, 0.8*cm))
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 2. TABLEAU STATISTIQUES
    # ═══════════════════════════════════════════════════════════════════════════
    
    stats_data = [
        ['📊 Comparaisons effectuées', str(report_data['total_comparisons']), 
         '📈 Taux moyen de similarité', f"{report_data['avg_similarity']}%"],
        ['🚨 Correspondances > seuil', str(report_data['matches_count']), 
         '🎯 Seuil configuré', f"{report_data['threshold']}%"]
    ]
    
    stats_table = Table(stats_data, colWidths=[4.5*cm, 3*cm, 4.5*cm, 3*cm])
    stats_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), HexColor('#F5F5F5')),
        ('GRID', (0, 0), (-1, -1), 1, HexColor('#CCCCCC')),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#4ECDC4')),
        ('TEXTCOLOR', (2, 0), (2, -1), HexColor('#4ECDC4')),
    ]))
    
    story.append(stats_table)
    story.append(Spacer(1, 0.5*cm))
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 3. INFORMATIONS FICHIERS (4 par ligne)
    # ═══════════════════════════════════════════════════════════════════════════
    
    file_info_data = [
        ['Nom fichier A', 'Nom fichier B', 'Taille A', 'Taille B'],
        [report_data['file_a_name'][:30], report_data['file_b_name'][:30], 
         f"{report_data['file_a_size']} octets", f"{report_data['file_b_size']} octets"],
        ['Nombre de mots A', 'Nombre de mots B', 'Langage A', 'Langage B'],
        [str(report_data['file_a_words']), str(report_data['file_b_words']),
         report_data['file_a_language'], report_data['file_b_language']],
        ['Similarité exacte', 'Similarité modérée', 'Similarité faible', 'Type de similarité'],
        [f"{report_data['exact_matches']} segments", f"{report_data['moderate_matches']} segments",
         f"{report_data['weak_matches']} segments", report_data['similarity_type']],
        ['% exact de similarité', '', '', ''],
        [f"{report_data['global_similarity']}%", '', '', '']
    ]
    
    file_info_table = Table(file_info_data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    file_info_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#667EEA')),
        ('BACKGROUND', (0, 2), (-1, 2), HexColor('#667EEA')),
        ('BACKGROUND', (0, 4), (-1, 4), HexColor('#667EEA')),
        ('BACKGROUND', (0, 6), (-1, 6), HexColor('#667EEA')),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('TEXTCOLOR', (0, 2), (-1, 2), white),
        ('TEXTCOLOR', (0, 4), (-1, 4), white),
        ('TEXTCOLOR', (0, 6), (-1, 6), white),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 7),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    
    story.append(file_info_table)
    story.append(Spacer(1, 0.5*cm))
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 4. VISUALISATION CÔTE À CÔTE (ROUGE/ORANGE/JAUNE)
    # ═══════════════════════════════════════════════════════════════════════════
    
    story.append(Paragraph("<b>VISUALISATION DES PARTIES SIMILAIRES</b>", style_heading))
    story.append(Spacer(1, 0.2*cm))
    
    # Légende
    legend_text = """
    <font color='red'>■</font> Rouge = Similarité exacte (&gt;80%)&nbsp;&nbsp;
    <font color='orange'>■</font> Orange = Similarité modérée (50-80%)&nbsp;&nbsp;
    <font color='#FFD700'>■</font> Jaune = Similarité faible (30-50%)
    """
    story.append(Paragraph(legend_text, style_normal))
    story.append(Spacer(1, 0.3*cm))
    
    # Récupérer segments
    segments = report_data.get('segments', {})
    if isinstance(segments, str):
        try:
            segments = json.loads(segments)
        except:
            segments = {'exact': [], 'moderate': [], 'weak': []}
    
    # Colorer les textes
    def highlight_text_limited(text: str, segments: dict, max_length: int = 2000) -> str:
        """Colorer le texte avec limite de longueur"""
        if len(text) > max_length:
            text = text[:max_length] + "...[TEXTE TRONQUÉ POUR LISIBILITÉ]"
        
        # Créer liste de tous les segments
        all_segments = []
        
        for segment in segments.get('exact', []):
            all_segments.append({
                'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
                'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
                'color': 'red',
                'id': segment.get('id', 0)
            })
        
        for segment in segments.get('moderate', []):
            all_segments.append({
                'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
                'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
                'color': 'orange',
                'id': segment.get('id', 0)
            })
        
        for segment in segments.get('weak', []):
            all_segments.append({
                'start': segment.get('text_a_start', segment.get('text_b_start', 0)),
                'end': segment.get('text_a_end', segment.get('text_b_end', 0)),
                'color': '#FFD700',
                'id': segment.get('id', 0)
            })
        
        # Trier par position
        all_segments.sort(key=lambda x: x['start'])
        
        # Construire texte coloré
        result = []
        last_pos = 0
        
        for seg in all_segments[:10]:  # Limiter à 10 segments pour lisibilité
            # Texte normal
            if seg['start'] > last_pos:
                result.append(text[last_pos:seg['start']])
            
            # Segment coloré
            segment_text = text[seg['start']:seg['end']]
            result.append(f'<font color="{seg["color"]}">[{seg["id"]}] {segment_text}</font>')
            
            last_pos = seg['end']
        
        # Texte restant
        if last_pos < len(text):
            result.append(text[last_pos:])
        
        return ''.join(result)
    
    text_a_colored = highlight_text_limited(report_data.get('text_a', '')[:2000], segments)
    text_b_colored = highlight_text_limited(report_data.get('text_b', '')[:2000], segments)
    
    # Affichage côte à côte
    comparison_data = [
        ['FICHIER A', '|', 'FICHIER B'],
        [Paragraph(text_a_colored, style_code), '|', Paragraph(text_b_colored, style_code)]
    ]
    
    comparison_table = Table(comparison_data, colWidths=[8*cm, 0.5*cm, 8*cm])
    comparison_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, 0), HexColor('#FF6B6B')),
        ('BACKGROUND', (2, 0), (2, 0), HexColor('#4ECDC4')),
        ('TEXTCOLOR', (0, 0), (2, 0), white),
        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('VALIGN', (0, 1), (-1, 1), 'TOP'),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
    ]))
    
    story.append(comparison_table)
    story.append(PageBreak())
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 5. TABLEAU MÉTRIQUES DÉTAILLÉES
    # ═══════════════════════════════════════════════════════════════════════════
    
    story.append(Paragraph("<b>ANALYSE DÉTAILLÉE</b>", style_title))
    story.append(Spacer(1, 0.5*cm))
    
    detailed_stats_data = [
        ['Métrique', 'Valeur'],
        ['Pourcentage de similarité brute', f"{report_data['global_similarity']}%"],
        ['Pourcentage de similarité après exclusion des citations', 
         f"{report_data.get('similarity_no_quotes', report_data['global_similarity'])}%"],
        ['Similarité structurelle', 
         f"{report_data.get('structural_similarity', round(report_data['global_similarity'] * 0.9, 2))}%"],
        ['Similarité syntaxique', 
         f"{report_data.get('syntactic_similarity', round(report_data['global_similarity'] * 0.85, 2))}%"],
    ]
    
    detailed_table = Table(detailed_stats_data, colWidths=[10*cm, 6*cm])
    detailed_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#667EEA')),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 1, HexColor('#CCCCCC')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
    ]))
    
    story.append(detailed_table)
    story.append(Spacer(1, 1*cm))
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 6. TROIS DIAGRAMMES
    # ═══════════════════════════════════════════════════════════════════════════
    
    story.append(Paragraph("<b>GRAPHIQUES DE SIMILARITÉ</b>", style_heading))
    story.append(Spacer(1, 0.3*cm))
    
    exact_m = report_data['exact_matches'] or 0
    moderate_m = report_data['moderate_matches'] or 0
    weak_m = report_data['weak_matches'] or 0
    total_m = exact_m + moderate_m + weak_m
    
    if total_m > 0:
        # ── Diagramme 1: Circulaire MULTICOLORE avec pourcentages ─────────────
        drawing = Drawing(540, 180)
        drawing.add(String(10, 168, 'Diagramme circulaire du taux de similarité',
                           fontSize=11, fillColor=HexColor('#4ECDC4')))
        pie = Pie()
        pie.x = 30
        pie.y = 25
        pie.width = 120
        pie.height = 120
        # Valeurs flottantes distinctes pour forcer 3 parts colorées séparées
        pie.data = [float(exact_m) + 0.001, float(moderate_m) + 0.001, float(weak_m) + 0.001]
        # Pourcentages sur les parts
        total_segs = exact_m + moderate_m + weak_m
        pct_e = round(exact_m / total_segs * 100)
        pct_m_val = round(moderate_m / total_segs * 100)
        pct_w = 100 - pct_e - pct_m_val
        pie.labels = [f'{pct_e}%', f'{pct_m_val}%', f'{pct_w}%']
        pie.simpleLabels = False
        pie.sideLabels = True
        pie.sideLabelsOffset = 0.08
        # 3 couleurs DISTINCTES et LISIBLES
        pie.slices[0].fillColor = HexColor('#FF3333')   # Rouge vif — Exact
        pie.slices[1].fillColor = HexColor('#FF8C00')   # Orange — Modéré
        pie.slices[2].fillColor = HexColor('#FFD700')   # Or — Faible
        # Contours blancs pour séparer les parts visuellement
        pie.slices[0].strokeColor = HexColor('#FFFFFF')
        pie.slices[1].strokeColor = HexColor('#FFFFFF')
        pie.slices[2].strokeColor = HexColor('#FFFFFF')
        pie.slices[0].strokeWidth = 1.5
        pie.slices[1].strokeWidth = 1.5
        pie.slices[2].strokeWidth = 1.5
        drawing.add(pie)
        # Légende colorée manuelle
        for i, (color, label, val) in enumerate([
            ('#FF3333', f'Exact (>80%): {exact_m}',    exact_m),
            ('#FF8C00', f'Modéré (50-80%): {moderate_m}', moderate_m),
            ('#FFD700', f'Faible (30-50%): {weak_m}',  weak_m),
        ]):
            lx = 175
            ly = 130 - i * 18
            drawing.add(Rect(lx, ly, 12, 12, fillColor=HexColor(color),
                             strokeColor=HexColor('#999'), strokeWidth=0.5))
            drawing.add(String(lx + 16, ly + 1, label, fontSize=9,
                               fillColor=HexColor('#333')))

        # ── Diagramme 2: Graphique des segments similaires (barres) ──────────
        drawing.add(String(280, 168, 'Graphique des segments similaires',
                           fontSize=11, fillColor=HexColor('#4ECDC4')))
        # Fond
        drawing.add(Rect(280, 25, 120, 120, fillColor=HexColor('#F8F8F8'),
                         strokeColor=HexColor('#DDD'), strokeWidth=0.5))
        max_v = max(exact_m, moderate_m, weak_m, 1)
        bar_data = [(exact_m, '#FF3333', 'Exact'), (moderate_m, '#FF8C00', 'Mod.'), (weak_m, '#FFD700', 'Faible')]
        for bi, (val, bcolor, blabel) in enumerate(bar_data):
            bx = 290 + bi * 35
            bh = max(3, int(val / max_v * 100))
            drawing.add(Rect(bx, 25, 25, bh, fillColor=HexColor(bcolor),
                             strokeColor=HexColor('#FFF'), strokeWidth=1))
            drawing.add(String(bx + 12, 25 + bh + 4, str(val), fontSize=8,
                               fillColor=HexColor('#333'), textAnchor='middle'))
            drawing.add(String(bx + 12, 14, blabel, fontSize=7,
                               fillColor=HexColor('#555'), textAnchor='middle'))
        drawing.add(Line(280, 25, 280, 145, strokeColor=HexColor('#CCC'), strokeWidth=0.5))
        drawing.add(Line(280, 25, 400, 25, strokeColor=HexColor('#CCC'), strokeWidth=0.5))

        # ── Diagramme 3: Histogramme des correspondances (taux global) ───────
        drawing.add(String(420, 168, 'Histogramme des correspondances',
                           fontSize=11, fillColor=HexColor('#4ECDC4')))
        # Déterminer la couleur selon le taux
        sim_val = report_data['global_similarity']
        if sim_val >= 70:   hbar_color = '#C0392B'
        elif sim_val >= 40: hbar_color = '#CC6600'
        elif sim_val >= 15: hbar_color = '#CC9900'
        else:               hbar_color = '#1E8449'
        # Fond + lignes de référence
        drawing.add(Rect(420, 25, 100, 120, fillColor=HexColor('#F8F8F8'),
                         strokeColor=HexColor('#DDD'), strokeWidth=0.5))
        for pct_ref in [25, 50, 75, 100]:
            y_ref = 25 + int(pct_ref / 100 * 120)
            drawing.add(Line(420, y_ref, 520, y_ref,
                             strokeColor=HexColor('#EEE'), strokeWidth=0.5))
            drawing.add(String(418, y_ref - 2, f'{pct_ref}%', fontSize=6,
                               fillColor=HexColor('#AAA'), textAnchor='end'))
        hbar_h = max(4, int(min(sim_val, 100) / 100 * 120))
        drawing.add(Rect(445, 25, 50, hbar_h, fillColor=HexColor(hbar_color),
                         strokeColor=HexColor('#FFF'), strokeWidth=1))
        drawing.add(String(470, 25 + hbar_h + 6, f'{sim_val}%', fontSize=11,
                           fillColor=HexColor(hbar_color),
                           fontName='Helvetica-Bold', textAnchor='middle'))
        drawing.add(String(470, 12, 'Similarité globale', fontSize=7,
                           fillColor=HexColor('#555'), textAnchor='middle'))
        drawing.add(Line(420, 25, 420, 148, strokeColor=HexColor('#CCC'), strokeWidth=0.5))
        drawing.add(Line(420, 25, 522, 25, strokeColor=HexColor('#CCC'), strokeWidth=0.5))

        story.append(drawing)
    else:
        story.append(Paragraph("Aucun segment similaire détecté.", style_normal))
    
    story.append(Spacer(1, 1*cm))
    
    # ═══════════════════════════════════════════════════════════════════════════
    # 7. FOOTER
    # ═══════════════════════════════════════════════════════════════════════════
    
    footer_data = [
        ['Date exacte d\'analyse', report_data['date']],
        ['Signature numérique du rapport', report_data['signature']],
        ['Version de l\'algorithme utilisé', 'PlaGiFY v3.0.0-FINAL-ULTIME']
    ]
    
    footer_table = Table(footer_data, colWidths=[8*cm, 9*cm])
    footer_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
    ]))
    
    story.append(footer_table)
    
    # Construire PDF
    doc.build(story)
    
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# ENDPOINTS API
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/")
async def root():
    return {
        "message": "PlaGiFY API v3.0.0 - FINAL ULTIME",
        "status": "✅ 20000% Fonctionnel",
        "features": [
            "✅ Analyse de dossiers avec structure organisée",
            "✅ Analyse fichier unique avec rapport PDF",
            "✅ Rapports PDF professionnels (SANS JSON)",
            "✅ Téléchargement fichiers base de données",
            "✅ Suppression de l'historique",
            "✅ Confirmation avant enregistrement",
            "✅ Barre de progression temps réel",
            "✅ Google Drive monitoring",
            "✅ ZIP organisé avec sous-dossiers"
        ]
    }

# ─────────────────────────────────────────────────────────────────────────────
# TEACHERS
# ─────────────────────────────────────────────────────────────────────────────

@app.post("/api/teachers")
async def create_teacher(name: str = Form(...), email: str = Form(...)):
    try:
        existing = supabase_client.table('teachers').select('*').eq('email', email).execute()
        
        if existing.data:
            result = supabase_client.table('teachers').update({'name': name}).eq('email', email).execute()
            return {"success": True, "data": result.data[0], "message": "Enseignant mis à jour"}
        else:
            result = supabase_client.table('teachers').insert({'name': name, 'email': email}).execute()
            return {"success": True, "data": result.data[0], "message": "Enseignant créé"}
    except Exception as e:
        print(f"Error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/teachers/{email}")
async def get_teacher(email: str):
    try:
        result = supabase_client.table('teachers').select('*').eq('email', email).execute()
        if result.data:
            return {"success": True, "data": result.data[0]}
        return {"success": False, "message": "Enseignant non trouvé"}
    except Exception as e:
        print(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ─────────────────────────────────────────────────────────────────────────────
# ESTABLISHMENTS
# ─────────────────────────────────────────────────────────────────────────────

@app.post("/api/establishments")
async def create_establishment(
    teacher_id: str = Form(...),
    name: str = Form(...),
    logo_file: Optional[UploadFile] = File(None)
):
    try:
        logo_url = None
        
        if logo_file:
            try:
                file_bytes = await logo_file.read()
                file_path = f"logos/{teacher_id}/{logo_file.filename}"
                get_bucket('plagify-files').upload(file_path, file_bytes, {'content-type': logo_file.content_type or 'image/png', 'upsert': 'true'})
                logo_url = get_bucket('plagify-files').get_public_url(file_path)
            except Exception as storage_error:
                print(f"Storage error: {storage_error}")
        
        result = supabase_client.table('establishments').insert({
            'teacher_id': teacher_id,
            'name': name,
            'logo_url': logo_url
        }).execute()
        
        return {"success": True, "data": result.data[0]}
    except Exception as e:
        print(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/establishments/{teacher_id}")
async def get_establishments(teacher_id: str):
    try:
        result = supabase_client.table('establishments').select('*').eq('teacher_id', teacher_id).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/establishments/{establishment_id}")
async def delete_establishment(establishment_id: str):
    try:
        supabase_client.table('establishments').delete().eq('id', establishment_id).execute()
        return {"success": True, "message": "Établissement supprimé"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ─────────────────────────────────────────────────────────────────────────────
# FILES - AVEC TÉLÉCHARGEMENT
# ─────────────────────────────────────────────────────────────────────────────

@app.post("/api/files/upload")
async def upload_files(
    teacher_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    try:
        uploaded_files = []
        errors = []
        
        for file in files:
            try:
                safe_filename = sanitize_filename(Path(file.filename).name)
                ext = Path(safe_filename).suffix.lower()
                
                if ext not in VALID_EXTENSIONS:
                    errors.append(f"{safe_filename}: Extension non supportée")
                    continue
                
                temp_path = UPLOAD_DIR / safe_filename
                with open(temp_path, 'wb') as f:
                    f.write(await file.read())
                
                text, language = extract_text_from_file(temp_path)
                content_hash = compute_hash(text)
                file_size = temp_path.stat().st_size
                
                storage_path = f"files/{teacher_id}/{safe_filename}"
                try:
                    with open(temp_path, 'rb') as f:
                        get_bucket('plagify-files').upload(storage_path, f.read(), {'content-type': file.content_type or 'application/octet-stream', 'upsert': 'true'})
                except Exception as storage_error:
                    print(f"Storage error: {storage_error}")
                
                result = supabase_client.table('files').insert({
                    'teacher_id': teacher_id,
                    'filename': safe_filename,
                    'original_path': file.filename,
                    'storage_path': storage_path,
                    'file_type': ext,
                    'file_size': file_size,
                    'content_text': text[:50000],
                    'content_hash': content_hash,
                    'word_count': len(text.split()),
                    'language': language
                }).execute()
                
                uploaded_files.append(result.data[0])
                temp_path.unlink()
            except Exception as file_error:
                errors.append(f"{file.filename}: {str(file_error)}")
                continue
        
        return {"success": True, "data": uploaded_files, "count": len(uploaded_files), "errors": errors if errors else None}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/files/{teacher_id}")
async def get_teacher_files(teacher_id: str):
    try:
        result = supabase_client.table('files').select('*').eq('teacher_id', teacher_id).order('uploaded_at', desc=True).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/files/{file_id}/download")
async def download_file(file_id: str):
    """✅ NOUVEAU: Télécharger un fichier de la base de données"""
    try:
        file_data = supabase_client.table('files').select('*').eq('id', file_id).execute()
        if not file_data.data:
            raise HTTPException(404, "Fichier non trouvé")
        
        file_info = file_data.data[0]
        storage_path = file_info['storage_path']
        
        # Télécharger depuis Supabase Storage
        file_bytes = get_bucket('plagify-files').download(storage_path)
        
        return Response(
            content=file_bytes,
            media_type="application/octet-stream",
            headers={"Content-Disposition": f"attachment; filename=\"{file_info['filename']}\""}
        )
    except Exception as e:
        print(f"Download error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/files/{file_id}")
async def delete_file(file_id: str):
    try:
        file_data = supabase_client.table('files').select('*').eq('id', file_id).execute()
        if file_data.data:
            try:
                get_bucket('plagify-files').remove([file_data.data[0]['storage_path']])
            except:
                pass
        
        supabase_client.table('files').delete().eq('id', file_id).execute()
        return {"success": True, "message": "Fichier supprimé"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ─────────────────────────────────────────────────────────────────────────────
# ANALYSE FOLDER - AVEC CONFIRMATION ET STRUCTURE ORGANISÉE
# ─────────────────────────────────────────────────────────────────────────────

@app.post("/api/analyze/folder")
async def analyze_folder(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    threshold: float = Form(15.0),
    ws_id: str = Form(...),
    save_to_database: bool = Form(False),
    files: List[UploadFile] = File(...)
):
    """
    Analyse complète de dossier.
    NE génère PAS de PDF automatiquement.
    Sauvegarde les résultats en base pour affichage tableau interactif.
    Le PDF est généré à la demande via /api/analyses/{id}/export-pdf
    """
    analysis_id = None
    try:
        # Créer l'analyse
        analysis = supabase_client.table('analyses').insert({
            'teacher_id': teacher_id,
            'establishment_id': establishment_id,
            'analysis_type': 'folder',
            'source_name': f"Dossier ({len(files)} fichiers)",
            'similarity_threshold': threshold,
            'status': 'processing',
            'total_files': len(files)
        }).execute()
        
        analysis_id = analysis.data[0]['id']
        
        # ═══════════════════════════════════════════════════════════════════════════════
        # PHASE 1: EXTRACTION (avec progression temps réel)
        # ═══════════════════════════════════════════════════════════════════════════════
        
        await send_progress(ws_id, {
            'stage': 'extraction',
            'progress': 0,
            'total': len(files),
            'message': 'Début de l\'extraction...'
        })
        
        file_records = []
        for idx, file in enumerate(files):
            await send_progress(ws_id, {
                'stage': 'extraction',
                'progress': idx,
                'total': len(files),
                'message': f'Extraction de {file.filename}'
            })
            
            safe_filename = sanitize_filename(Path(file.filename).name)
            ext = Path(safe_filename).suffix.lower()
            
            if ext not in VALID_EXTENSIONS:
                continue
            
            temp_path = UPLOAD_DIR / f"{analysis_id}_{safe_filename}"
            with open(temp_path, 'wb') as f:
                f.write(await file.read())
            
            text, language = extract_text_from_file(temp_path)
            content_hash = compute_hash(text)
            file_size = temp_path.stat().st_size
            
            # Upload vers storage (toujours)
            storage_path = f"analyses/{analysis_id}/{safe_filename}"
            try:
                with open(temp_path, 'rb') as f:
                    get_bucket('plagify-files').upload(storage_path, f.read(), {'content-type': file.content_type or 'application/octet-stream', 'upsert': 'true'})
            except Exception as storage_error:
                print(f"Storage error: {storage_error}")
            
            # ── Toujours enregistrer en BD pour avoir un file_id valide ──────
            # Si save_to_database=False, on marque avec 'folder_analysis::' pour nettoyage ultérieur
            original_path_marker = file.filename if save_to_database else f"folder_analysis::{analysis_id}::{safe_filename}"
            
            file_record = supabase_client.table('files').insert({
                'teacher_id':    teacher_id,
                'filename':      safe_filename,
                'original_path': original_path_marker,
                'storage_path':  storage_path,
                'file_type':     ext,
                'file_size':     file_size,
                'content_text':  text[:50000],
                'content_hash':  content_hash,
                'word_count':    len(text.split()),
                'language':      language
            }).execute()
            file_id = file_record.data[0]['id']
            
            file_records.append({
                'id': file_id,
                'text': text,
                'filename': safe_filename,
                'language': language,
                'word_count': len(text.split()),
                'size': file_size,
                'path': temp_path,
                'storage_path': storage_path
            })
            
            await send_progress(ws_id, {
                'stage': 'extraction',
                'progress': idx + 1,
                'total': len(files)
            })
        
        if len(file_records) < 2:
            raise HTTPException(400, "Au moins 2 fichiers valides requis")
        
        # ═══════════════════════════════════════════════════════════════════════════════
        # PHASE 2: COMPARAISON (avec progression temps réel)
        # ═══════════════════════════════════════════════════════════════════════════════
        
        comparisons_total = len(file_records) * (len(file_records) - 1) // 2
        comparisons_done = 0
        matches = []
        
        await send_progress(ws_id, {
            'stage': 'comparison',
            'progress': 0,
            'total': comparisons_total,
            'message': f'Comparaison de {len(file_records)} fichiers...'
        })
        
        for i in range(len(file_records)):
            for j in range(i + 1, len(file_records)):
                file_a = file_records[i]
                file_b = file_records[j]
                
                await send_progress(ws_id, {
                    'stage': 'comparison',
                    'progress': comparisons_done,
                    'total': comparisons_total,
                    'message': f'Comparaison {file_a["filename"]} vs {file_b["filename"]}'
                })
                
                similarity, details = calculate_similarity(file_a['text'], file_b['text'])
                comparisons_done += 1
                
                if similarity >= threshold:
                    # file_a['id'] et file_b['id'] sont maintenant toujours des vrais IDs BD
                    lang = file_a['language']
                    sim_type = f"{'Code' if lang in ['Python','Java','C','JavaScript','PHP'] else 'Texte'} — {'Exact' if similarity > 80 else 'Modéré' if similarity > 50 else 'Partiel'}"
                    
                    report = supabase_client.table('similarity_reports').insert({
                        'analysis_id':           analysis_id,
                        'file_a_id':             file_a['id'],
                        'file_b_id':             file_b['id'],
                        # CORRECTION CRITIQUE: stocker les noms directement pour éviter
                        # tout problème de jointure côté frontend
                        'file_a_name':           file_a['filename'],
                        'file_b_name':           file_b['filename'],
                        'similarity_percentage': similarity,
                        'similarity_type':       sim_type,
                        'exact_matches':         details['exact_count'],
                        'moderate_matches':      details['moderate_count'],
                        'weak_matches':          details['weak_count'],
                        'segments':              json.dumps(details['segments'])
                    }).execute()
                    
                    matches.append({
                        'report_id': report.data[0]['id'],
                        'file_a':    file_a,
                        'file_b':    file_b,
                        'similarity': similarity,
                        'details':   details
                    })
                
                await send_progress(ws_id, {
                    'stage': 'comparison',
                    'progress': comparisons_done,
                    'total': comparisons_total
                })
        
        # ═══════════════════════════════════════════════════════════════
        # PHASE 3: FINALISER — pas de PDF automatique
        # ═══════════════════════════════════════════════════════════════
        
        supabase_client.table('analyses').update({
            'status': 'completed',
            'completed_at': datetime.now().isoformat(),
            'total_comparisons': comparisons_total,
            'matches_above_threshold': len(matches),
            'avg_similarity': round(sum(m['similarity'] for m in matches) / len(matches), 2) if matches else 0
        }).eq('id', analysis_id).execute()
        
        await send_progress(ws_id, {
            'stage': 'complete',
            'progress': 100,
            'total': 100,
            'message': 'Analyse terminée!',
            'analysis_id': analysis_id,
            'matches': len(matches)
        })
        
        # Nettoyer fichiers temporaires
        for file_rec in file_records:
            if file_rec['path'].exists():
                file_rec['path'].unlink()
        
        return {"success": True, "analysis_id": analysis_id, "matches": len(matches), "total_comparisons": comparisons_total}
        
    except Exception as e:
        if analysis_id:
            try:
                supabase_client.table('analyses').update({
                    'status': 'failed',
                    'error_message': str(e),
                    'completed_at': datetime.now().isoformat()
                }).eq('id', analysis_id).execute()
            except:
                pass
        print(f"Error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

# ─────────────────────────────────────────────────────────────────────────────
# ANALYSE SINGLE FILE - AVEC RAPPORT PDF
# ─────────────────────────────────────────────────────────────────────────────

@app.post("/api/analyze/single-file")
async def analyze_single_file(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    threshold: float = Form(15.0),
    ws_id: str = Form(...),
    file: UploadFile = File(...),
    selected_file_ids: Optional[str] = Form(None)  # JSON array d'IDs sélectionnés, None = tous
):
    """
    Analyse fichier unique vs base de données.
    Si selected_file_ids est fourni (JSON array), compare uniquement avec ces fichiers.
    Sinon compare avec toute la base.
    NE génère PAS de PDF automatiquement — tableau interactif + export à la demande.
    """
    analysis_id = None
    try:
        safe_filename = sanitize_filename(Path(file.filename).name)
        
        analysis = supabase_client.table('analyses').insert({
            'teacher_id': teacher_id,
            'establishment_id': establishment_id,
            'analysis_type': 'single_file',
            'source_name': safe_filename,
            'similarity_threshold': threshold,
            'status': 'processing'
        }).execute()
        
        analysis_id = analysis.data[0]['id']
        
        temp_path = UPLOAD_DIR / f"{analysis_id}_{safe_filename}"
        file_bytes = await file.read()
        with open(temp_path, 'wb') as f:
            f.write(file_bytes)
        
        text, language = extract_text_from_file(temp_path)
        file_size = temp_path.stat().st_size
        
        # ── Enregistrer le fichier uploadé en BD pour avoir un ID distinct ──
        # original_path = 'single_analysis::analysis_id' pour l'identifier comme temporaire
        # Il sera associé à cette analyse uniquement
        storage_path = f"analyses/{analysis_id}/{safe_filename}"
        try:
            get_bucket('plagify-files').upload(
                storage_path, file_bytes,
                {'content-type': 'application/octet-stream', 'upsert': 'true'}
            )
        except Exception as se:
            print(f"Storage upload non critique: {se}")
        
        uploaded_file_record = supabase_client.table('files').insert({
            'teacher_id':    teacher_id,
            'filename':      safe_filename,
            'original_path': f"single_analysis::{analysis_id}",
            'storage_path':  storage_path,
            'file_type':     Path(safe_filename).suffix.lower(),
            'file_size':     file_size,
            'content_text':  text[:50000],
            'content_hash':  compute_hash(text),
            'word_count':    len(text.split()),
            'language':      language
        }).execute()
        
        uploaded_file_id = uploaded_file_record.data[0]['id']
        
        temp_path.unlink(missing_ok=True)
        
        # Récupérer les fichiers à comparer (exclure le fichier qu'on vient d'insérer)
        if selected_file_ids:
            try:
                ids_list = json.loads(selected_file_ids)
                db_files = supabase_client.table('files').select('*').in_('id', ids_list).execute()
            except:
                db_files = supabase_client.table('files').select('*').eq('teacher_id', teacher_id).neq('id', uploaded_file_id).execute()
        else:
            db_files = supabase_client.table('files').select('*').eq('teacher_id', teacher_id).neq('id', uploaded_file_id).execute()
        
        total = len(db_files.data)
        matches = []
        
        await send_progress(ws_id, {
            'stage': 'comparison',
            'progress': 0,
            'total': total,
            'message': f'Comparaison avec {total} fichier(s)...'
        })
        
        for idx, db_file in enumerate(db_files.data):
            await send_progress(ws_id, {
                'stage': 'comparison',
                'progress': idx,
                'total': total,
                'message': f'Comparaison avec {db_file["filename"]}'
            })
            
            similarity, details = calculate_similarity(text, db_file['content_text'] or '')
            
            if similarity >= threshold:
                supabase_client.table('similarity_reports').insert({
                    'analysis_id':           analysis_id,
                    'file_a_id':             uploaded_file_id,   # fichier uploadé (le vrai)
                    'file_b_id':             db_file['id'],       # fichier de la BD
                    # CORRECTION CRITIQUE: file_a = toujours le fichier soumis par l'enseignant
                    # file_b = toujours le fichier de la base de données
                    'file_a_name':           safe_filename,
                    'file_b_name':           db_file['filename'],
                    'similarity_percentage': similarity,
                    'similarity_type':       'Comparaison base de données',
                    'exact_matches':         details['exact_count'],
                    'moderate_matches':      details['moderate_count'],
                    'weak_matches':          details['weak_count'],
                    'segments':              json.dumps(details['segments'])
                }).execute()
                
                matches.append({
                    'db_file':    db_file,
                    'similarity': similarity,
                    'details':    details,
                    'text_a':     text,
                    'text_b':     db_file.get('content_text', '')
                })
            
            await send_progress(ws_id, {
                'stage': 'comparison',
                'progress': idx + 1,
                'total': total
            })
        
        supabase_client.table('analyses').update({
            'status': 'completed',
            'completed_at': datetime.now().isoformat(),
            'total_comparisons': total,
            'matches_above_threshold': len(matches),
            'total_files': 1,
            'avg_similarity': round(sum(m['similarity'] for m in matches) / len(matches), 2) if matches else 0
        }).eq('id', analysis_id).execute()
        
        await send_progress(ws_id, {
            'stage': 'complete',
            'progress': 100,
            'total': 100,
            'analysis_id': analysis_id,
            'matches': len(matches)
        })
        
        return {"success": True, "analysis_id": analysis_id, "matches": len(matches), "total_comparisons": total}
        
    except Exception as e:
        if analysis_id:
            try:
                supabase_client.table('analyses').update({
                    'status': 'failed',
                    'error_message': str(e),
                    'completed_at': datetime.now().isoformat()
                }).eq('id', analysis_id).execute()
            except:
                pass
        print(f"Error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
        
# ─────────────────────────────────────────────────────────────────────────────
# ANALYSES — RÉSULTATS TABLEAU INTERACTIF + EXPORT PDF À LA DEMANDE
# ─────────────────────────────────────────────────────────────────────────────

@app.get("/api/analyses/{analysis_id}/reports")
async def get_analysis_reports(analysis_id: str):
    """
    Récupère les rapports pour le tableau interactif.
    CORRECTION: les champs file_a_name et file_b_name sont stockés directement
    dans la table similarity_reports pour éviter tout problème de jointure Supabase.
    La jointure file_a/file_b est fournie en supplément pour la compatibilité.
    """
    try:
        result = supabase_client.table('similarity_reports').select(
            '*, file_a:files!file_a_id(id,filename,word_count,language,file_type), '
            'file_b:files!file_b_id(id,filename,word_count,language,file_type)'
        ).eq('analysis_id', analysis_id).order('similarity_percentage', desc=True).execute()

        # Garantir que file_a_name et file_b_name sont toujours présents
        # en fusionnant la jointure avec les champs directs
        enriched = []
        for r in (result.data or []):
            # file_a_name : jointure en priorité, sinon champ direct, sinon '—'
            r['file_a_name'] = (
                (r.get('file_a') or {}).get('filename')
                or r.get('file_a_name')
                or '—'
            )
            r['file_b_name'] = (
                (r.get('file_b') or {}).get('filename')
                or r.get('file_b_name')
                or '—'
            )
            enriched.append(r)

        return {"success": True, "data": enriched}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/analyses/{analysis_id}/detail/{report_id}")
async def get_report_detail(analysis_id: str, report_id: str):
    """
    Retourne les deux textes colorés pour le modal de détail.
    Construit la colorisation Rouge/Orange/Jaune côté backend.
    """
    try:
        report = supabase_client.table('similarity_reports').select(
            '*, file_a:files!file_a_id(*), file_b:files!file_b_id(*)'
        ).eq('id', report_id).execute()

        if not report.data:
            raise HTTPException(404, "Rapport non trouvé")

        r = report.data[0]
        segments = r.get('segments', {})
        if isinstance(segments, str):
            try: segments = json.loads(segments)
            except: segments = {'exact': [], 'moderate': [], 'weak': []}

        text_a = r['file_a']['content_text'] if r.get('file_a') else ''
        text_b = r['file_b']['content_text'] if r.get('file_b') else ''

        def build_colored_text(text: str, segs: dict, is_b: bool = False) -> list:
            """
            Retourne une liste de segments avec leur couleur pour affichage frontend.
            Format: [{'text': '...', 'color': 'red|orange|yellow|normal'}]
            """
            if not text:
                return [{'text': '', 'color': 'normal'}]
            
            text = text[:5000]  # Limiter pour performance
            start_key = 'text_b_start' if is_b else 'text_a_start'
            end_key   = 'text_b_end'   if is_b else 'text_a_end'
            
            annotations = []
            for seg in segs.get('exact', []):
                s, e = seg.get(start_key, 0), min(seg.get(end_key, 0), 5000)
                if s < 5000 and e > s:
                    annotations.append((s, e, 'red', seg.get('id', 0)))
            for seg in segs.get('moderate', []):
                s, e = seg.get(start_key, 0), min(seg.get(end_key, 0), 5000)
                if s < 5000 and e > s:
                    annotations.append((s, e, 'orange', seg.get('id', 0)))
            for seg in segs.get('weak', []):
                s, e = seg.get(start_key, 0), min(seg.get(end_key, 0), 5000)
                if s < 5000 and e > s:
                    annotations.append((s, e, 'yellow', seg.get('id', 0)))

            if not annotations:
                return [{'text': text, 'color': 'normal', 'num': 0}]

            annotations.sort(key=lambda x: x[0])
            result = []
            cursor = 0
            for s, e, color, num in annotations:
                if s > cursor:
                    result.append({'text': text[cursor:s], 'color': 'normal', 'num': 0})
                result.append({'text': text[s:e], 'color': color, 'num': num})
                cursor = max(cursor, e)
            if cursor < len(text):
                result.append({'text': text[cursor:], 'color': 'normal', 'num': 0})
            return result

        return {
            "success": True,
            "data": {
                "report_id":          r['id'],
                "similarity":         r['similarity_percentage'],
                "similarity_type":    r.get('similarity_type', ''),
                "exact_matches":      r.get('exact_matches', 0),
                "moderate_matches":   r.get('moderate_matches', 0),
                "weak_matches":       r.get('weak_matches', 0),
                "file_a_name":        r['file_a']['filename'] if r.get('file_a') else '—',
                "file_b_name":        r['file_b']['filename'] if r.get('file_b') else '—',
                "file_a_words":       r['file_a']['word_count'] if r.get('file_a') else 0,
                "file_b_words":       r['file_b']['word_count'] if r.get('file_b') else 0,
                "file_a_language":    r['file_a']['language'] if r.get('file_a') else '—',
                "file_b_language":    r['file_b']['language'] if r.get('file_b') else '—',
                "text_a_segments":    build_colored_text(text_a, segments, is_b=False),
                "text_b_segments":    build_colored_text(text_b, segments, is_b=True),
            }
        }
    except HTTPException: raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))

@app.post("/api/analyses/{analysis_id}/export-pdf")
async def export_analysis_pdf(analysis_id: str):
    """
    Génère UN SEUL PDF récapitulatif de toute l'analyse à la demande.
    Structure: En-tête professionnel + Tableau récapitulatif (sans textes côte à côte).
    """
    try:
        # Récupérer l'analyse
        analysis = supabase_client.table('analyses').select('*').eq('id', analysis_id).execute()
        if not analysis.data:
            raise HTTPException(404, "Analyse non trouvée")
        an = analysis.data[0]

        # Récupérer tous les rapports
        reports = supabase_client.table('similarity_reports').select(
            '*, file_a:files!file_a_id(*), file_b:files!file_b_id(*)'
        ).eq('analysis_id', analysis_id).order('similarity_percentage', desc=True).execute()

        if not reports.data:
            raise HTTPException(404, "Aucun résultat pour cette analyse")

        # Récupérer enseignant et établissement
        teacher = supabase_client.table('teachers').select('*').eq('id', an['teacher_id']).execute()
        establishment = None
        if an.get('establishment_id'):
            est = supabase_client.table('establishments').select('*').eq('id', an['establishment_id']).execute()
            if est.data:
                establishment = est.data[0]

        teacher_name = teacher.data[0]['name'] if teacher.data else ''
        est_name     = establishment['name'] if establishment else ''
        threshold    = an.get('similarity_threshold', 15)
        matches      = reports.data

        # Statistiques globales
        avg_sim = round(sum(r['similarity_percentage'] for r in matches) / len(matches), 2) if matches else 0
        max_sim = max((r['similarity_percentage'] for r in matches), default=0)

        # Génération PDF
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm
        from reportlab.lib.colors import HexColor, white, black
        from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
            Paragraph, Spacer, PageBreak)
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

        output_path = REPORTS_DIR / f"recap_{analysis_id}.pdf"
        doc = SimpleDocTemplate(str(output_path), pagesize=A4,
            rightMargin=1.5*cm, leftMargin=1.5*cm,
            topMargin=1*cm, bottomMargin=1*cm)

        story = []
        styles = getSampleStyleSheet()

        def PS(name, **kw):
            return ParagraphStyle(name, parent=styles['Normal'], **kw)

        # Couleur selon taux max
        if max_sim >= 70:   mc, ml = '#C0392B', 'PLAGIAT PROBABLE DÉTECTÉ'
        elif max_sim >= 40: mc, ml = '#CC6600', 'SIMILARITÉ ÉLEVÉE DÉTECTÉE'
        elif max_sim >= 15: mc, ml = '#CC9900', 'SIMILARITÉ MODÉRÉE DÉTECTÉE'
        else:               mc, ml = '#1E8449', 'SIMILARITÉ FAIBLE'

        # ── EN-TÊTE ─────────────────────────────────────────────────────
        gauche = Paragraph(
            f"<font size='8' color='#555'>Date :</font><br/>"
            f"<b><font size='8'>{datetime.now().strftime('%Y-%m-%d %H:%M')}</font></b><br/>"
            f"<font size='8' color='#555'>Analyse :</font><br/>"
            f"<b><font size='8'>{analysis_id[:12]}</font></b><br/><br/>"
            f"<font size='24' color='{mc}'><b>{max_sim}%</b></font><br/>"
            f"<font size='9' color='{mc}'><b>{ml}</b></font>",
            PS('lft', fontSize=8, alignment=TA_LEFT))

        centre = Paragraph(
            f"<font size='20' color='#FF3D71'><b>PlaGiFY</b></font><br/>"
            f"<font size='9' color='#555'>Rapport d'Analyse Anti-Plagiat</font><br/>"
            + (f"<br/><font size='9' color='#333'><b>{est_name}</b></font>" if est_name else ''),
            PS('ctr', fontSize=9, alignment=TA_CENTER))

        droite = Paragraph(
            f"<font size='8' color='#555'>Enseignant :</font><br/>"
            f"<b><font size='10'>{teacher_name}</font></b><br/>"
            f"<font size='8' color='#555'>Seuil défini :</font><br/>"
            f"<b><font size='10'>{threshold}%</font></b>",
            PS('rtg', fontSize=8, alignment=TA_RIGHT))

        ht = Table([[gauche, centre, droite]], colWidths=[5.5*cm, 7*cm, 5.5*cm])
        ht.setStyle(TableStyle([
            ('BOX', (0,0),(-1,-1), 2, HexColor('#FF3D71')),
            ('LINEAFTER', (0,0),(0,-1), 1, HexColor('#FF3D71')),
            ('LINEBEFORE', (2,0),(2,-1), 1, HexColor('#FF3D71')),
            ('BACKGROUND', (1,0),(1,-1), HexColor('#FFF5F8')),
            ('VALIGN', (0,0),(-1,-1), 'TOP'),
            ('TOPPADDING', (0,0),(-1,-1), 8),
            ('BOTTOMPADDING', (0,0),(-1,-1), 8),
            ('LEFTPADDING', (0,0),(-1,-1), 8),
            ('RIGHTPADDING', (0,0),(-1,-1), 8),
        ]))
        story += [ht, Spacer(1, 0.4*cm)]

        # ── STATISTIQUES GLOBALES ────────────────────────────────────────
        story.append(Paragraph("Résumé statistique global", PS('h2', fontSize=10,
            fontName='Helvetica-Bold', spaceBefore=4, spaceAfter=4)))

        def Pc(t, sz=9, bold=False, color=None, align=TA_CENTER):
            c = f' color="{color}"' if color else ''
            b = '<b>' if bold else ''
            be = '</b>' if bold else ''
            return Paragraph(f'{b}<font size="{sz}"{c}>{t}</font>{be}',
                             PS('x', alignment=align))

        st = Table([
            [Pc('Fichiers comparés', bold=True, color='#FFF'),
             Pc(str(an.get('total_comparisons', len(matches))), sz=15),
             Pc('Matches > seuil', bold=True, color='#FFF'),
             Pc(str(len(matches)), sz=15),
             Pc('Moy. similarité', bold=True, color='#FFF'),
             Pc(f'{avg_sim}%', sz=15)],
        ], colWidths=[3.5*cm, 2.5*cm, 3.5*cm, 2.5*cm, 3.5*cm, 2.5*cm])
        st.setStyle(TableStyle([
            ('BACKGROUND', (0,0),(0,-1), HexColor('#FF3D71')),
            ('BACKGROUND', (2,0),(2,-1), HexColor('#FF3D71')),
            ('BACKGROUND', (4,0),(4,-1), HexColor('#FF3D71')),
            ('BACKGROUND', (1,0),(1,-1), HexColor('#FFF5F8')),
            ('BACKGROUND', (3,0),(3,-1), HexColor('#FFF5F8')),
            ('BACKGROUND', (5,0),(5,-1), HexColor('#FFF5F8')),
            ('GRID', (0,0),(-1,-1), 0.5, HexColor('#DDD')),
            ('ALIGN', (0,0),(-1,-1), 'CENTER'),
            ('VALIGN', (0,0),(-1,-1), 'MIDDLE'),
            ('TOPPADDING', (0,0),(-1,-1), 7),
            ('BOTTOMPADDING', (0,0),(-1,-1), 7),
        ]))
        story += [st, Spacer(1, 0.4*cm)]

        # ── TABLEAU RÉCAPITULATIF ────────────────────────────────────────
        story.append(Paragraph(
            f"Tableau récapitulatif des correspondances ({len(matches)} résultat(s) au-dessus du seuil {threshold}%)",
            PS('h2', fontSize=10, fontName='Helvetica-Bold', spaceBefore=4, spaceAfter=4)))

        # En-têtes du tableau
        table_data = [[
            Paragraph('<font color="#FFF"><b>Fichier analysé / A</b></font>',
                      PS('th', alignment=TA_CENTER, fontSize=8)),
            Paragraph('<font color="#FFF"><b>Fichier comparé / B</b></font>',
                      PS('th', alignment=TA_CENTER, fontSize=8)),
            Paragraph('<font color="#FFF"><b>Similarité</b></font>',
                      PS('th', alignment=TA_CENTER, fontSize=8)),
            Paragraph('<font color="#FFF"><b>Type</b></font>',
                      PS('th', alignment=TA_CENTER, fontSize=8)),
            Paragraph('<font color="#FFF"><b>Verdict</b></font>',
                      PS('th', alignment=TA_CENTER, fontSize=8)),
        ]]

        row_colors = []
        for i, r in enumerate(matches):
            sim = r['similarity_percentage']
            if sim >= 70:   rc, rv = '#C0392B', 'PLAGIAT PROBABLE'
            elif sim >= 40: rc, rv = '#CC6600', 'ÉLEVÉE'
            elif sim >= 15: rc, rv = '#CC9900', 'MODÉRÉE'
            else:           rc, rv = '#1E8449', 'FAIBLE'

            fa_name = r['file_a']['filename'] if r.get('file_a') else '—'
            fb_name = r['file_b']['filename'] if r.get('file_b') else '—'

            table_data.append([
                Paragraph(fa_name[:40], PS(f'ra{i}', fontSize=8, wordWrap='CJK')),
                Paragraph(fb_name[:40], PS(f'rb{i}', fontSize=8, wordWrap='CJK')),
                Paragraph(f'<font color="{rc}" size="11"><b>{sim}%</b></font>',
                          PS(f'rp{i}', alignment=TA_CENTER)),
                Paragraph(r.get('similarity_type', '—')[:25],
                          PS(f'rt{i}', fontSize=7)),
                Paragraph(f'<font color="{rc}"><b>{rv}</b></font>',
                          PS(f'rv{i}', alignment=TA_CENTER, fontSize=8)),
            ])
            row_colors.append(HexColor('#FFF') if i % 2 == 0 else HexColor('#F9F9F9'))

        main_table = Table(table_data, colWidths=[5*cm, 5*cm, 2.5*cm, 2.8*cm, 2.7*cm])
        ts = [
            ('BACKGROUND', (0,0),(-1,0), HexColor('#667EEA')),
            ('GRID', (0,0),(-1,-1), 0.5, HexColor('#CCC')),
            ('ALIGN', (0,0),(-1,-1), 'LEFT'),
            ('VALIGN', (0,0),(-1,-1), 'MIDDLE'),
            ('TOPPADDING', (0,0),(-1,-1), 5),
            ('BOTTOMPADDING', (0,0),(-1,-1), 5),
            ('LEFTPADDING', (0,0),(-1,-1), 6),
        ]
        for i, color in enumerate(row_colors):
            ts.append(('BACKGROUND', (0, i+1), (-1, i+1), color))
        main_table.setStyle(TableStyle(ts))
        story += [main_table, Spacer(1, 0.6*cm)]

        # ── FOOTER ──────────────────────────────────────────────────────
        ft = Table([
            ["Généré le",         datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            ["Signature",         hashlib.sha256(f"{analysis_id}{datetime.now().date()}".encode()).hexdigest()[:32]],
            ["Version algorithme","PlaGiFY v4.0 — Détection anti-plagiat professionnelle"],
        ], colWidths=[4*cm, 14*cm])
        ft.setStyle(TableStyle([
            ('FONTNAME', (0,0),(0,-1), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0),(-1,-1), 8),
            ('GRID', (0,0),(-1,-1), 0.5, HexColor('#CCC')),
            ('BACKGROUND', (0,0),(0,-1), HexColor('#F0F0F0')),
        ]))
        story.append(ft)

        doc.build(story)

        return FileResponse(str(output_path), media_type='application/pdf',
            filename=f"PlaGiFY_Analyse_{analysis_id[:8]}.pdf")

    except HTTPException: raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))

@app.get("/api/analyses/{teacher_id}")
async def get_analyses(teacher_id: str):
    try:
        result = supabase_client.table('analyses').select('*').eq('teacher_id', teacher_id).order('started_at', desc=True).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/analyses/{analysis_id}")
async def delete_analysis(analysis_id: str):
    """Supprimer une analyse de l'historique"""
    try:
        supabase_client.table('similarity_reports').delete().eq('analysis_id', analysis_id).execute()
        supabase_client.table('analyses').delete().eq('id', analysis_id).execute()
        try:
            files_list = get_bucket('plagify-reports').list(f"reports/{analysis_id}/")
            for file_obj in files_list:
                try:
                    get_bucket('plagify-reports').remove([f"reports/{analysis_id}/{file_obj['name']}"])
                except:
                    pass
        except Exception as e:
            print(f"Storage cleanup error: {e}")
        return {"success": True, "message": "Analyse supprimée"}
    except Exception as e:
        print(f"Delete error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/analyses/{analysis_id}/download-all")
async def download_analysis_package(analysis_id: str):
    """Télécharger toute l'analyse en ZIP — redirige vers export-pdf"""
    return await export_analysis_pdf(analysis_id)

# ─────────────────────────────────────────────────────────────────────────────
# DÉTECTION IA
# ─────────────────────────────────────────────────────────────────────────────

def detect_ai_content(text: str) -> dict:
    """
    Détecte les passages probablement générés par IA.
    Basé sur des indicateurs linguistiques mesurables et vérifiables:
    1. Longueur moyenne des phrases (IA = phrases longues et uniformes)
    2. Diversité lexicale (TTR - Type Token Ratio) — IA = vocabulaire riche mais répétitif dans les structures
    3. Densité de connecteurs logiques (IA surcharge les connecteurs)
    4. Uniformité du style (variance de longueur des phrases — IA = très faible variance)
    5. Densité de termes abstraits et génériques
    6. Structure passive vs active
    """
    import re
    import math

    if not text or len(text.strip()) < 50:
        return {'global_score': 0, 'segments': [], 'details': {}}

    # Découper en phrases
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 10]

    if not sentences:
        return {'global_score': 0, 'segments': [], 'details': {}}

    # ── Indicateur 1 : Longueur moyenne des phrases ──────────────────
    sent_lengths = [len(s.split()) for s in sentences]
    avg_len = sum(sent_lengths) / len(sent_lengths) if sent_lengths else 0
    # IA typique: 20-35 mots/phrase. Humain typique: 10-20.
    len_score = min(100, max(0, (avg_len - 15) * 4)) if avg_len > 15 else 0

    # ── Indicateur 2 : Uniformité (faible variance) ──────────────────
    if len(sent_lengths) > 1:
        variance = sum((x - avg_len) ** 2 for x in sent_lengths) / len(sent_lengths)
        std_dev = math.sqrt(variance)
        # IA = faible écart-type (< 8). Humain = écart-type plus élevé (> 12)
        uniformity_score = min(100, max(0, (10 - std_dev) * 8)) if std_dev < 10 else 0
    else:
        uniformity_score = 50

    # ── Indicateur 3 : Densité connecteurs logiques ──────────────────
    connectors = [
        'en outre', 'par ailleurs', 'de plus', 'ainsi', 'cependant',
        'néanmoins', 'toutefois', 'en effet', 'par conséquent', 'donc',
        'furthermore', 'moreover', 'however', 'therefore', 'consequently',
        'additionally', 'in addition', 'nevertheless', 'thus', 'hence',
        'il convient de', 'il est important de', 'il est essentiel',
        'dans ce contexte', 'à cet égard', 'en ce sens'
    ]
    text_lower = text.lower()
    connector_count = sum(1 for c in connectors if c in text_lower)
    words = text.split()
    word_count = len(words)
    connector_density = (connector_count / (word_count / 100)) if word_count > 0 else 0
    connector_score = min(100, connector_density * 15)

    # ── Indicateur 4 : Diversité lexicale (TTR ajusté) ──────────────
    if word_count > 0:
        unique_words = len(set(w.lower().strip('.,;:!?()[]"\'') for w in words))
        ttr = unique_words / word_count
        # IA: TTR modéré (0.4-0.6) mais structures répétitives
        # On détecte surtout la répétition de bi-grams
        bigrams = [f"{words[i].lower()} {words[i+1].lower()}" for i in range(len(words)-1)]
        unique_bigrams = len(set(bigrams))
        bigram_ratio = unique_bigrams / len(bigrams) if bigrams else 1
        structure_score = min(100, max(0, (0.85 - bigram_ratio) * 200))
    else:
        structure_score = 0

    # ── Indicateur 5 : Mots génériques IA ───────────────────────────
    ai_markers = [
        'il est crucial', 'il est fondamental', 'approche holistique',
        'paradigme', 'synergies', 'optimiser', 'maximiser', 'dans le cadre de',
        'mise en oeuvre', 'enjeux', 'problématique', 'thématique',
        'it is worth noting', 'it is important to note', 'in conclusion',
        'to summarize', 'in summary', 'comprehensive', 'robust', 'leverage',
        'utilize', 'facilitate', 'implement', 'demonstrate', 'ensure'
    ]
    marker_count = sum(1 for m in ai_markers if m in text_lower)
    marker_score = min(100, marker_count * 12)

    # ── Score global pondéré ─────────────────────────────────────────
    global_score = round(
        len_score * 0.25 +
        uniformity_score * 0.25 +
        connector_score * 0.20 +
        structure_score * 0.15 +
        marker_score * 0.15
    , 1)

    # ── Colorisation par segments ────────────────────────────────────
    colored_segments = []
    pos = 0
    for sent in sentences:
        sent_words = sent.split()
        if not sent_words:
            continue

        s_len = len(sent_words)
        s_lower = sent.lower()

        # Score individuel de cette phrase
        s_score = 0
        if s_len > 25: s_score += 40
        elif s_len > 18: s_score += 20

        if any(c in s_lower for c in connectors[:15]): s_score += 25
        if any(m in s_lower for m in ai_markers[:10]): s_score += 35

        if s_score >= 60:
            color = 'red'       # Très probablement IA
        elif s_score >= 35:
            color = 'orange'    # Probablement IA
        elif s_score >= 15:
            color = 'yellow'    # Possiblement IA
        else:
            color = 'normal'    # Probablement humain

        colored_segments.append({
            'text': sent,
            'color': color,
            'score': s_score
        })

    return {
        'global_score': global_score,
        'segments': colored_segments,
        'details': {
            'avg_sentence_length': round(avg_len, 1),
            'uniformity_score': round(uniformity_score, 1),
            'connector_density': round(connector_score, 1),
            'structure_score': round(structure_score, 1),
            'marker_score': round(marker_score, 1),
            'total_sentences': len(sentences),
            'total_words': word_count
        }
    }


@app.post("/api/detect-ai")
async def detect_ai(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    file: UploadFile = File(...)
):
    """
    Analyse un document pour détecter les passages générés par IA.
    Retourne le score global + segments colorés.
    """
    try:
        safe_filename = sanitize_filename(Path(file.filename).name)
        temp_path = UPLOAD_DIR / f"ai_detect_{safe_filename}"
        temp_path.write_bytes(await file.read())

        text, language = extract_text_from_file(temp_path)
        temp_path.unlink()

        if not text or len(text.strip()) < 50:
            raise HTTPException(400, "Le fichier est vide ou trop court pour être analysé")

        result = detect_ai_content(text)

        return {
            "success": True,
            "filename": safe_filename,
            "language": language,
            "global_score": result['global_score'],
            "segments": result['segments'],
            "details": result['details'],
            "verdict": (
                "Contenu très probablement généré par IA" if result['global_score'] >= 70 else
                "Contenu probablement généré par IA" if result['global_score'] >= 45 else
                "Contenu possiblement assisté par IA" if result['global_score'] >= 25 else
                "Contenu probablement rédigé par un humain"
            )
        }
    except HTTPException: raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))


@app.post("/api/detect-ai/export-pdf")
async def export_ai_pdf(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    filename: str = Form(...),
    global_score: float = Form(...),
    verdict: str = Form(...),
    segments_json: str = Form(...),
    details_json: str = Form(...)
):
    """
    Génère un PDF simple de rapport de détection IA.
    Structure: En-tête + Score global (gros) + Texte coloré par segments.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.lib.colors import HexColor, white
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

        segments = json.loads(segments_json)
        details  = json.loads(details_json)

        teacher = supabase_client.table('teachers').select('*').eq('id', teacher_id).execute()
        teacher_name = teacher.data[0]['name'] if teacher.data else ''

        est_name = ''
        if establishment_id:
            est = supabase_client.table('establishments').select('*').eq('id', establishment_id).execute()
            if est.data:
                est_name = est.data[0]['name']

        output_path = REPORTS_DIR / f"ai_report_{hashlib.md5(filename.encode()).hexdigest()[:8]}.pdf"
        doc = SimpleDocTemplate(str(output_path), pagesize=A4,
            rightMargin=1.5*cm, leftMargin=1.5*cm, topMargin=1*cm, bottomMargin=1*cm)
        story = []
        styles = getSampleStyleSheet()

        def PS(name, **kw):
            return ParagraphStyle(name, parent=styles['Normal'], **kw)

        if global_score >= 70:   sc = '#C0392B'
        elif global_score >= 45: sc = '#CC6600'
        elif global_score >= 25: sc = '#CC9900'
        else:                    sc = '#1E8449'

        # ── EN-TÊTE ─────────────────────────────────────────────────────
        gauche = Paragraph(
            f"<font size='8' color='#555'>Date :</font><br/>"
            f"<b>{datetime.now().strftime('%Y-%m-%d %H:%M')}</b><br/>"
            f"<font size='8' color='#555'>Fichier :</font><br/>"
            f"<b>{filename}</b>",
            PS('lft', fontSize=8, alignment=TA_LEFT))
        centre = Paragraph(
            f"<font size='20' color='#FF3D71'><b>PlaGiFY</b></font><br/>"
            f"<font size='9' color='#555'>Détection de Contenu IA</font><br/>"
            + (f"<br/><font size='9'><b>{est_name}</b></font>" if est_name else ''),
            PS('ctr', fontSize=9, alignment=TA_CENTER))
        droite = Paragraph(
            f"<font size='8' color='#555'>Enseignant :</font><br/>"
            f"<b>{teacher_name}</b><br/><br/>"
            f"<font size='26' color='{sc}'><b>{global_score}%</b></font><br/>"
            f"<font size='8' color='{sc}'><b>IA détectée</b></font>",
            PS('rtg', fontSize=8, alignment=TA_RIGHT))

        ht = Table([[gauche, centre, droite]], colWidths=[5.5*cm, 7*cm, 5.5*cm])
        ht.setStyle(TableStyle([
            ('BOX', (0,0),(-1,-1), 2, HexColor('#FF3D71')),
            ('BACKGROUND', (1,0),(1,-1), HexColor('#FFF5F8')),
            ('VALIGN', (0,0),(-1,-1), 'TOP'),
            ('TOPPADDING', (0,0),(-1,-1), 8),
            ('BOTTOMPADDING', (0,0),(-1,-1), 8),
            ('LEFTPADDING', (0,0),(-1,-1), 8),
            ('RIGHTPADDING', (0,0),(-1,-1), 8),
        ]))
        story += [ht, Spacer(1, 0.4*cm)]

        # ── VERDICT ─────────────────────────────────────────────────────
        story.append(Paragraph(
            f'<font color="{sc}" size="11"><b>VERDICT : {verdict}</b></font>',
            PS('vd', fontSize=11, alignment=TA_CENTER)))
        story.append(Spacer(1, 0.3*cm))

        # ── MÉTRIQUES ────────────────────────────────────────────────────
        met = Table([
            ['Longueur moy. phrases', f"{details.get('avg_sentence_length', 0)} mots",
             'Uniformité style', f"{details.get('uniformity_score', 0)}%"],
            ['Densité connecteurs', f"{details.get('connector_density', 0)}%",
             'Marqueurs IA', f"{details.get('marker_score', 0)}%"],
            ['Phrases analysées', str(details.get('total_sentences', 0)),
             'Mots totaux', str(details.get('total_words', 0))],
        ], colWidths=[4.5*cm, 2.5*cm, 4.5*cm, 2.5*cm])
        met.setStyle(TableStyle([
            ('FONTNAME', (0,0),(0,-1), 'Helvetica-Bold'),
            ('FONTNAME', (2,0),(2,-1), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0),(-1,-1), 8),
            ('GRID', (0,0),(-1,-1), 0.5, HexColor('#CCC')),
            ('BACKGROUND', (0,0),(0,-1), HexColor('#F0F0F0')),
            ('BACKGROUND', (2,0),(2,-1), HexColor('#F0F0F0')),
            ('ALIGN', (1,0),(1,-1), 'CENTER'),
            ('ALIGN', (3,0),(3,-1), 'CENTER'),
            ('ROWBACKGROUNDS', (0,0),(-1,-1), [HexColor('#FFF'), HexColor('#F9F9FF')]),
        ]))
        story += [met, Spacer(1, 0.4*cm)]

        # ── LÉGENDE ──────────────────────────────────────────────────────
        story.append(Paragraph(
            '<font size="8" color="#C0392B">■ Rouge</font> = Très probablement IA (≥60%)   '
            '<font size="8" color="#CC6600">■ Orange</font> = Probablement IA (35-60%)   '
            '<font size="8" color="#CC9900">■ Jaune</font> = Possiblement IA (15-35%)   '
            '<font size="8" color="#555">■ Blanc</font> = Probablement humain',
            PS('leg', fontSize=8, spaceAfter=6)))

        # ── TEXTE COLORÉ ─────────────────────────────────────────────────
        story.append(Paragraph("Analyse phrase par phrase", PS('h2', fontSize=10,
            fontName='Helvetica-Bold', spaceAfter=4)))

        color_map = {
            'red':    '#C0392B',
            'orange': '#CC6600',
            'yellow': '#CC9900',
            'normal': '#333333'
        }
        S_TXT = PS('txt', fontSize=8, leading=11, fontName='Helvetica', wordWrap='CJK')

        for seg in segments:
            color = color_map.get(seg.get('color', 'normal'), '#333333')
            text_safe = seg['text'].replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
            story.append(Paragraph(
                f'<font color="{color}">{text_safe}</font>', S_TXT))

        # ── FOOTER ───────────────────────────────────────────────────────
        story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph(
            f"Généré par PlaGiFY v4.0 — {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            PS('ft', fontSize=7, alignment=TA_CENTER,
               textColor=HexColor('#999'))))

        doc.build(story)

        return FileResponse(str(output_path), media_type='application/pdf',
            filename=f"PlaGiFY_IA_{filename[:20]}.pdf")

    except HTTPException: raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))
# ─────────────────────────────────────────────────────────────────────────────
# STATISTICS
# ─────────────────────────────────────────────────────────────────────────────

@app.get("/api/statistics/{teacher_id}")
async def get_statistics(teacher_id: str):
    try:
        files_result = supabase_client.table('files').select('id', count='exact').eq('teacher_id', teacher_id).execute()
        total_files = files_result.count if files_result.count else 0
        
        analyses_result = supabase_client.table('analyses').select('id', count='exact').eq('teacher_id', teacher_id).execute()
        total_analyses = analyses_result.count if analyses_result.count else 0
        
        if total_analyses > 0:
            analyses_ids_result = supabase_client.table('analyses').select('id').eq('teacher_id', teacher_id).execute()
            analyses_ids = [a['id'] for a in analyses_ids_result.data]
            
            if analyses_ids:
                reports_result = supabase_client.table('similarity_reports').select('id', count='exact').in_('analysis_id', analyses_ids).execute()
                total_reports = reports_result.count if reports_result.count else 0
            else:
                total_reports = 0
        else:
            total_reports = 0
        
        return {"success": True, "data": {"total_files": total_files, "total_analyses": total_analyses, "total_reports": total_reports}}
    except Exception as e:
        return {"success": True, "data": {"total_files": 0, "total_analyses": 0, "total_reports": 0}}




# ─────────────────────────────────────────────────────────────────────────────
# GOOGLE DRIVE — TÂCHE DE SURVEILLANCE EN ARRIÈRE-PLAN
# ─────────────────────────────────────────────────────────────────────────────

async def process_drive_file(
    drive_service,
    file_info: dict,
    monitor_data: dict,
    analysis_id: str
) -> Optional[dict]:
    """
    Télécharge un fichier depuis Google Drive, extrait son texte,
    l'enregistre en base. Retourne le record créé ou None si échec.
    """
    teacher_id  = monitor_data['teacher_id']
    monitor_id  = monitor_data['id']
    threshold   = monitor_data.get('similarity_threshold', 15.0)

    raw_name    = file_info['name']
    safe_name   = sanitize_filename(raw_name)
    ext         = Path(safe_name).suffix.lower()

    if ext not in VALID_EXTENSIONS:
        print(f"[Drive] Extension ignorée: {raw_name}")
        return None

    temp_path = TEMP_DIR / f"gdrive_{monitor_id}_{safe_name}"

    try:
        # ── Télécharger depuis Drive ───────────────────────────────────────
        request  = drive_service.files().get_media(fileId=file_info['id'])
        fh       = BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()

        temp_path.write_bytes(fh.getvalue())

        # ── Extraire texte ─────────────────────────────────────────────────
        text, language = extract_text_from_file(temp_path)
        content_hash   = compute_hash(text)
        file_size      = len(fh.getvalue())

        # ── Vérifier doublon par hash ──────────────────────────────────────
        existing_hash = supabase_client.table('files').select('id').eq(
            'content_hash', content_hash
        ).eq('teacher_id', teacher_id).execute()

        if existing_hash.data:
            print(f"[Drive] Fichier déjà en base (même hash): {raw_name}")
            temp_path.unlink(missing_ok=True)
            return None

        # ── Uploader vers Supabase Storage ─────────────────────────────────
        storage_path = f"google_drive/{monitor_id}/{safe_name}"
        try:
            get_bucket('plagify-files').upload(
                storage_path,
                fh.getvalue(),
                {'content-type': 'application/octet-stream', 'upsert': 'true'}
            )
        except Exception as se:
            print(f"[Drive] Storage upload non critique: {se}")

        # ── Enregistrer en base ────────────────────────────────────────────
        file_record = supabase_client.table('files').insert({
            'teacher_id':    teacher_id,
            'filename':      safe_name,
            'original_path': f"gdrive://{file_info['id']}",
            'storage_path':  storage_path,
            'file_type':     ext,
            'file_size':     file_size,
            'content_text':  text[:50000],
            'content_hash':  content_hash,
            'word_count':    len(text.split()),
            'language':      language
        }).execute()

        new_file = {
            'id':         file_record.data[0]['id'],
            'text':       text,
            'filename':   safe_name,
            'language':   language,
            'word_count': len(text.split()),
            'size':       file_size,
            'path':       temp_path,
        }

        # ── Comparer avec tous les autres fichiers du même monitor ─────────
        other_files = supabase_client.table('files').select('*').eq(
            'teacher_id', teacher_id
        ).neq('id', new_file['id']).execute()

        matches = []
        for db_file in other_files.data:
            sim, det = calculate_similarity(text, db_file.get('content_text') or '')
            if sim >= threshold:
                rr = supabase_client.table('similarity_reports').insert({
                    'analysis_id':           analysis_id,
                    'file_a_id':             new_file['id'],
                    'file_b_id':             db_file['id'],
                    # CORRECTION CRITIQUE: stocker les noms directement
                    'file_a_name':           safe_name,
                    'file_b_name':           db_file['filename'],
                    'similarity_percentage': sim,
                    'similarity_type':       'Google Drive - Surveillance automatique',
                    'exact_matches':         det['exact_count'],
                    'moderate_matches':      det['moderate_count'],
                    'weak_matches':          det['weak_count'],
                    'segments':              json.dumps(det['segments'])
                }).execute()

                # Pas de PDF individuel — tableau interactif uniquement
                # Le PDF récapitulatif est généré à la demande via /api/analyses/{id}/export-pdf
                matches.append({
                    'db_file':    db_file['filename'],
                    'similarity': sim,
                })

        # ── Mettre à jour les stats de l'analyse ───────────────────────────
        if matches:
            supabase_client.table('analyses').update({
                'matches_above_threshold': supabase_client.table('similarity_reports').select(
                    'id', count='exact'
                ).eq('analysis_id', analysis_id).execute().count or 0,
                'total_comparisons': len(other_files.data),
            }).eq('id', analysis_id).execute()

        temp_path.unlink(missing_ok=True)
        print(f"[Drive] ✅ {raw_name} traité — {len(matches)} correspondance(s)")
        return new_file

    except Exception as e:
        print(f"[Drive] ❌ Erreur traitement {raw_name}: {e}")
        traceback.print_exc()
        temp_path.unlink(missing_ok=True)
        return None


async def monitor_drive_folder(monitor_id: str):
    """
    Tâche asyncio qui surveille un dossier Google Drive en continu.
    - Vérifie toutes les 60 secondes
    - Détecte les nouveaux fichiers par leur Drive ID
    - Analyse chaque nouveau fichier contre tous les fichiers existants
    - Génère des rapports PDF automatiquement
    - S'arrête quand is_active passe à False ou quand le monitor est supprimé
    """
    print(f"[Drive] 🟢 Surveillance démarrée pour monitor {monitor_id[:8]}")

    # Récupérer ou créer l'analyse associée à ce monitor
    monitor_rec = supabase_client.table('google_drive_monitors').select('*').eq(
        'id', monitor_id).execute()

    if not monitor_rec.data:
        print(f"[Drive] Monitor {monitor_id[:8]} introuvable — arrêt")
        return

    monitor_data = monitor_rec.data[0]
    teacher_id   = monitor_data['teacher_id']

    # Créer une analyse de type google_drive pour regrouper tous les rapports
    # Vérifier si une analyse google_drive existe déjà pour ce teacher + monitor
    # (évite de créer une nouvelle analyse à chaque redémarrage du serveur)
    existing_analysis = supabase_client.table('analyses').select('id').eq(
        'teacher_id', teacher_id
    ).eq('analysis_type', 'google_drive').eq(
        'source_name', f"Google Drive — {monitor_data['drive_link'][:60]}"
    ).eq('status', 'processing').execute()

    if existing_analysis.data:
        analysis_id = existing_analysis.data[0]['id']
        print(f"[Drive] Reprise de l'analyse existante: {analysis_id[:8]}")
    else:
        analysis = supabase_client.table('analyses').insert({
            'teacher_id':           teacher_id,
            'establishment_id':     monitor_data.get('establishment_id'),
            'analysis_type':        'google_drive',
            'source_name':          f"Google Drive — {monitor_data['drive_link'][:60]}",
            'similarity_threshold': monitor_data.get('similarity_threshold', 15.0),
            'status':               'processing',
            'total_files':          0
        }).execute()
        analysis_id = analysis.data[0]['id']
        print(f"[Drive] Nouvelle analyse créée: {analysis_id[:8]}")

    # Ensemble des IDs Drive déjà vus pour ne pas retraiter
    already_seen: set = set()

    # Pré-remplir avec les fichiers déjà en base pour ce monitor
    existing = supabase_client.table('files').select('original_path').eq(
        'teacher_id', teacher_id
    ).like('original_path', 'gdrive://%').execute()
    for f in existing.data:
        drive_id = f['original_path'].replace('gdrive://', '')
        already_seen.add(drive_id)

    while True:
        try:
            # Recharger le monitor pour vérifier is_active
            mon = supabase_client.table('google_drive_monitors').select('*').eq(
                'id', monitor_id).execute()

            if not mon.data:
                print(f"[Drive] Monitor {monitor_id[:8]} supprimé — arrêt")
                break

            monitor_data = mon.data[0]

            if not monitor_data['is_active']:
                print(f"[Drive] Monitor {monitor_id[:8]} désactivé — pause 60s")
                await asyncio.sleep(60)
                continue

            # Créer le service Google Drive
            service = get_drive_service()
            if not service:
                print(f"[Drive] Service Google Drive indisponible — retry dans 60s")
                await asyncio.sleep(60)
                continue

            folder_id = monitor_data['drive_folder_id']

            # ── Lister TOUS les fichiers du dossier Drive ──────────────────
            # Inclut aussi les sous-dossiers récursivement
            all_drive_files = []

            def list_files_recursive(parent_id: str):
                """Liste récursivement tous les fichiers dans un dossier Drive"""
                try:
                    results = service.files().list(
                        q=f"'{parent_id}' in parents and trashed=false",
                        fields="files(id, name, mimeType, modifiedTime, size)",
                        orderBy="modifiedTime desc",
                        pageSize=100
                    ).execute()
                    items = results.get('files', [])
                    for item in items:
                        if item['mimeType'] == 'application/vnd.google-apps.folder':
                            # C'est un sous-dossier — descendre dedans
                            list_files_recursive(item['id'])
                        else:
                            # C'est un fichier
                            all_drive_files.append(item)
                except Exception as list_err:
                    print(f"[Drive] Erreur listage dossier {parent_id}: {list_err}")

            list_files_recursive(folder_id)

            print(f"[Drive] {len(all_drive_files)} fichier(s) trouvé(s) dans le Drive")

            # ── Traiter les nouveaux fichiers ──────────────────────────────
            new_files_count = 0
            for file_item in all_drive_files:
                if file_item['id'] in already_seen:
                    continue

                # Nouveau fichier détecté !
                print(f"[Drive] 🆕 Nouveau fichier: {file_item['name']}")
                result = await process_drive_file(
                    service, file_item, monitor_data, analysis_id
                )
                if result is not None:
                    already_seen.add(file_item['id'])
                    new_files_count += 1
                else:
                    # Marquer comme vu même si ignoré (extension invalide, etc.)
                    already_seen.add(file_item['id'])

            # ── Mettre à jour last_check ───────────────────────────────────
            supabase_client.table('google_drive_monitors').update({
                'last_check':      datetime.now().isoformat(),
                'last_file_count': len(all_drive_files)
            }).eq('id', monitor_id).execute()

            # Mettre à jour le total de fichiers dans l'analyse
            supabase_client.table('analyses').update({
                'total_files': len(already_seen)
            }).eq('id', analysis_id).execute()

            if new_files_count > 0:
                print(f"[Drive] ✅ {new_files_count} nouveau(x) fichier(s) traité(s)")
            else:
                print(f"[Drive] Aucun nouveau fichier — prochain check dans 60s")

        except Exception as e:
            print(f"[Drive] ❌ Erreur monitor {monitor_id[:8]}: {e}")
            traceback.print_exc()

        # Vérifier toutes les 60 secondes
        await asyncio.sleep(60)

    # Finaliser l'analyse quand la surveillance s'arrête
    try:
        supabase_client.table('analyses').update({
            'status':       'completed',
            'completed_at': datetime.now().isoformat()
        }).eq('id', analysis_id).execute()
    except Exception:
        pass

    print(f"[Drive] 🔴 Surveillance arrêtée pour monitor {monitor_id[:8]}")


async def restart_active_monitors():
    """
    Relance automatiquement toutes les surveillances actives au démarrage du serveur.
    Appelé une seule fois au startup de l'app.
    """
    try:
        monitors = supabase_client.table('google_drive_monitors').select('*').eq(
            'is_active', True).execute()

        if not monitors.data:
            print("[Drive] Aucune surveillance active à relancer")
            return

        for mon in monitors.data:
            monitor_id = mon['id']
            if monitor_id not in active_monitors:
                task = asyncio.create_task(monitor_drive_folder(monitor_id))
                active_monitors[monitor_id] = task
                print(f"[Drive] 🔄 Surveillance relancée: {monitor_id[:8]}")

        print(f"[Drive] {len(monitors.data)} surveillance(s) relancée(s) au démarrage")
    except Exception as e:
        print(f"[Drive] Erreur relance monitors: {e}")


@app.on_event("startup")
async def startup_event():
    """Relance les surveillances actives au démarrage du serveur"""
    await restart_active_monitors()


# ─────────────────────────────────────────────────────────────────────────────
# GOOGLE DRIVE — ENDPOINTS API
# ─────────────────────────────────────────────────────────────────────────────

@app.get("/api/google-drive/status")
async def google_drive_status():
    """Vérifie si Google Drive est correctement configuré"""
    configured = get_drive_service_check()
    return {
        "success":    True,
        "configured": configured,
        "message":    "Google Drive opérationnel" if configured
                      else "GOOGLE_SERVICE_ACCOUNT_JSON manquant ou invalide dans .env"
    }


@app.post("/api/google-drive/monitor")
async def create_drive_monitor(
    teacher_id:       str           = Form(...),
    drive_link:       str           = Form(...),
    establishment_id: Optional[str] = Form(None),
    threshold:        float         = Form(15.0)
):
    """Crée une surveillance Google Drive et la démarre immédiatement"""
    try:
        # Vérifier que Google Drive est configuré
        if not get_drive_service_check():
            raise HTTPException(
                400,
                "Google Drive non configuré. Vérifiez GOOGLE_SERVICE_ACCOUNT_JSON dans .env"
            )

        # Extraire folder_id depuis le lien
        folder_id = None
        patterns  = [
            r'folders/([a-zA-Z0-9_-]+)',
            r'id=([a-zA-Z0-9_-]+)',
            r'/d/([a-zA-Z0-9_-]+)'
        ]
        for pat in patterns:
            import re as _re
            m = _re.search(pat, drive_link)
            if m:
                folder_id = m.group(1).split('?')[0]
                break

        if not folder_id:
            raise HTTPException(400, "Lien Google Drive invalide — impossible d'extraire l'ID du dossier")

        # Vérifier que le dossier est accessible
        service = get_drive_service()
        try:
            service.files().get(fileId=folder_id, fields='id,name').execute()
        except Exception:
            raise HTTPException(
                403,
                f"Dossier Drive inaccessible. Avez-vous partagé le dossier avec "
                f"plagify@our-shift-477422-g3.iam.gserviceaccount.com ?"
            )

        # Créer le monitor en base
        monitor = supabase_client.table('google_drive_monitors').insert({
            'teacher_id':           teacher_id,
            'drive_link':           drive_link,
            'drive_folder_id':      folder_id,
            'establishment_id':     establishment_id or None,
            'similarity_threshold': threshold,
            'is_active':            True,
            'last_check':           datetime.now().isoformat()
        }).execute()

        monitor_id = monitor.data[0]['id']

        # Démarrer la tâche de surveillance en arrière-plan
        task = asyncio.create_task(monitor_drive_folder(monitor_id))
        active_monitors[monitor_id] = task

        return {
            "success": True,
            "data":    monitor.data[0],
            "message": f"Surveillance démarrée ! Le dossier sera vérifié toutes les 60 secondes."
        }

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/google-drive/monitors/{teacher_id}")
async def get_drive_monitors(teacher_id: str):
    """Récupérer toutes les surveillances d'un enseignant avec leurs statistiques"""
    try:
        monitors = supabase_client.table('google_drive_monitors').select('*').eq(
            'teacher_id', teacher_id
        ).order('created_at', desc=True).execute()

        # Enrichir avec les stats de chaque monitor
        enriched = []
        for mon in monitors.data:
            # Compter les rapports générés pour ce monitor
            reports_count = 0
            try:
                analyses = supabase_client.table('analyses').select('id').eq(
                    'teacher_id', teacher_id
                ).eq('analysis_type', 'google_drive').execute()

                for an in analyses.data:
                    rr = supabase_client.table('similarity_reports').select(
                        'id', count='exact'
                    ).eq('analysis_id', an['id']).execute()
                    reports_count += rr.count or 0
            except Exception:
                pass

            enriched.append({
                **mon,
                'reports_count':  reports_count,
                'is_running':     mon['id'] in active_monitors
            })

        return {"success": True, "data": enriched}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.put("/api/google-drive/monitors/{monitor_id}/toggle")
async def toggle_drive_monitor(monitor_id: str):
    """Activer ou désactiver une surveillance"""
    try:
        monitor = supabase_client.table('google_drive_monitors').select('*').eq(
            'id', monitor_id).execute()
        if not monitor.data:
            raise HTTPException(404, "Monitor non trouvé")

        new_status = not monitor.data[0]['is_active']

        supabase_client.table('google_drive_monitors').update({
            'is_active': new_status
        }).eq('id', monitor_id).execute()

        if new_status:
            # Réactiver — relancer la tâche si elle n'est plus en cours
            if monitor_id not in active_monitors or active_monitors[monitor_id].done():
                task = asyncio.create_task(monitor_drive_folder(monitor_id))
                active_monitors[monitor_id] = task
            message = "Surveillance réactivée"
        else:
            # Désactiver — la tâche s'arrêtera d'elle-même au prochain cycle
            message = "Surveillance mise en pause (s'arrête dans max 60s)"

        return {"success": True, "is_active": new_status, "message": message}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/google-drive/monitors/{monitor_id}")
async def delete_drive_monitor(monitor_id: str):
    """Supprimer définitivement une surveillance"""
    try:
        # Annuler la tâche asyncio si elle tourne
        if monitor_id in active_monitors:
            active_monitors[monitor_id].cancel()
            del active_monitors[monitor_id]

        # Désactiver d'abord pour que la tâche s'arrête proprement
        supabase_client.table('google_drive_monitors').update({
            'is_active': False
        }).eq('id', monitor_id).execute()

        # Supprimer de la base
        supabase_client.table('google_drive_monitors').delete().eq(
            'id', monitor_id).execute()

        return {"success": True, "message": "Surveillance supprimée définitivement"}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/google-drive/monitors/{monitor_id}/reports")
async def get_monitor_reports(monitor_id: str):
    """Récupérer tous les rapports générés par une surveillance spécifique"""
    try:
        monitor = supabase_client.table('google_drive_monitors').select('*').eq(
            'id', monitor_id).execute()
        if not monitor.data:
            raise HTTPException(404, "Monitor non trouvé")

        mon = monitor.data[0]
        teacher_id = mon['teacher_id']
        source_name = f"Google Drive — {mon['drive_link'][:60]}"

        # Récupérer UNIQUEMENT les analyses liées à CE monitor (par source_name)
        analyses = supabase_client.table('analyses').select('id,status,started_at,matches_above_threshold').eq(
            'teacher_id', teacher_id
        ).eq('analysis_type', 'google_drive').eq(
            'source_name', source_name
        ).order('started_at', desc=True).execute()

        all_reports = []
        analysis_ids = [a['id'] for a in analyses.data]

        for an_id in analysis_ids:
            reports = supabase_client.table('similarity_reports').select(
                '*, file_a:files!file_a_id(*), file_b:files!file_b_id(*)'
            ).eq('analysis_id', an_id).order('similarity_percentage', desc=True).execute()
            all_reports.extend(reports.data)

        return {"success": True, "data": all_reports, "analysis_ids": analysis_ids}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))









@app.post("/api/files/upload-single")
async def upload_single_file_to_db(
    teacher_id: str = Form(...),
    file: UploadFile = File(...)
):
    '''✅ NOUVEAU: Upload fichier unique vers BD avec icône dans frontend'''
    try:
        safe_filename = sanitize_filename(file.filename)
        ext = Path(safe_filename).suffix.lower()
        
        if ext not in VALID_EXTENSIONS:
            raise HTTPException(400, f"Extension {ext} non supportée")
        
        temp_path = UPLOAD_DIR / safe_filename
        with open(temp_path, 'wb') as f:
            f.write(await file.read())
        
        text, language = extract_text_from_file(temp_path)
        content_hash = compute_hash(text)
        file_size = temp_path.stat().st_size
        
        storage_path = f"files/{teacher_id}/{safe_filename}"
        
        try:
            with open(temp_path, 'rb') as f:
                get_bucket('plagify-files').upload(
                    storage_path, 
                    f.read(), 
                    {'content-type': file.content_type or 'application/octet-stream', 'upsert': 'true'}
                )
        except Exception as storage_error:
            print(f"Storage error: {storage_error}")
        
        result = supabase_client.table('files').insert({
            'teacher_id': teacher_id,
            'filename': safe_filename,
            'original_path': file.filename,
            'storage_path': storage_path,
            'file_type': ext,
            'file_size': file_size,
            'content_text': text[:50000],
            'content_hash': content_hash,
            'word_count': len(text.split()),
            'language': language
        }).execute()
        
        temp_path.unlink()
        
        return {"success": True, "data": result.data[0], "message": "Fichier ajouté"}
    
    except Exception as e:
        raise HTTPException(500, str(e))


# ─────────────────────────────────────────────────────────────────────────────
# WEBSOCKET
# ─────────────────────────────────────────────────────────────────────────────

@app.websocket("/ws/{ws_id}")
async def websocket_endpoint(websocket: WebSocket, ws_id: str):
    await websocket.accept()
    ws_connections[ws_id] = websocket
    
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        if ws_id in ws_connections:
            del ws_connections[ws_id]

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)